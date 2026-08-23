"""PPT/PPTX text extraction using python-pptx.

Extracts the title and all text-box/placeholder/bullet text from every
slide, in shape order (title first when present). Slides whose text is
embedded only inside a picture are not OCR'd in this MVP (documented as
a known limitation) — a slide with genuinely no extractable text is
simply skipped.

Note: python-pptx supports .pptx (the modern, XML-based, zip format).
Legacy binary .ppt files are not readable by python-pptx; such files
raise DocumentProcessingError with a clear message asking the user to
save the file as .pptx.
"""
from __future__ import annotations

from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from app.core.exceptions import DocumentProcessingError
from app.core.logging import get_logger
from app.ingestion.base import ExtractedUnit
from app.ingestion.normalization import is_effectively_empty, normalize_text

logger = get_logger(__name__)


def extract_pptx(file_path: Path) -> List[ExtractedUnit]:
    if file_path.suffix.lower() == ".ppt":
        raise DocumentProcessingError(
            "Legacy binary .ppt files are not supported. Please re-save the "
            "file as .pptx (File > Save As > PowerPoint Presentation) and "
            "upload again."
        )

    try:
        presentation = Presentation(str(file_path))
    except PackageNotFoundError as exc:
        raise DocumentProcessingError(f"Invalid or corrupt PPTX file: {exc}") from exc
    except Exception as exc:
        raise DocumentProcessingError(f"Could not open PPTX: {exc}") from exc

    units: List[ExtractedUnit] = []
    for slide_index, slide in enumerate(presentation.slides):
        slide_number = slide_index + 1
        parts: List[str] = []

        title_text = _extract_title(slide)
        if title_text:
            parts.append(title_text)

        for shape in slide.shapes:
            if shape.has_text_frame and shape != getattr(
                slide.shapes, "title", None
            ):
                shape_text = _extract_shape_text(shape)
                if shape_text and shape_text != title_text:
                    parts.append(shape_text)
            if shape.has_table:
                parts.append(_extract_table_text(shape))

        raw_text = "\n".join(p for p in parts if p)
        normalized = normalize_text(raw_text)
        if is_effectively_empty(normalized):
            continue

        units.append(ExtractedUnit(text=normalized, slide_number=slide_number))

    return units


def _extract_title(slide) -> str:
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            return slide.shapes.title.text_frame.text.strip()
    except Exception:  # pragma: no cover - defensive
        pass
    return ""


def _extract_shape_text(shape) -> str:
    lines = []
    for paragraph in shape.text_frame.paragraphs:
        line = "".join(run.text for run in paragraph.runs) or paragraph.text
        if line.strip():
            lines.append(line.strip())
    return "\n".join(lines)


def _extract_table_text(shape) -> str:
    rows_text = []
    for row in shape.table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        rows_text.append(" | ".join(c for c in cells if c))
    return "\n".join(r for r in rows_text if r)
